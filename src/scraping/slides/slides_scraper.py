import os
import json
import time
import hashlib
from pathlib import Path
from typing import List

import requests
from dotenv import load_dotenv
from pptx import Presentation
from pdfminer.high_level import extract_text

# --------------------------------------------------
# CONFIG
# --------------------------------------------------
load_dotenv()

SERPAPI_KEY = os.getenv("SERPAPI_KEY")
if not SERPAPI_KEY:
    raise RuntimeError("SERPAPI_KEY missing")

OUTPUT_DIR = Path("scraping_output/slides_files")
OUTPUT_JSONL = Path("scraping_output/slides_dataset.jsonl")
STATE_FILE = Path("scraping_output/slides_state.json")

OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_JSONL.parent.mkdir(parents=True, exist_ok=True)

TARGET_SLIDES = 300
RESULTS_PER_QUERY = 10
SLEEP = 2

QUERIES = [
    "machine learning lecture filetype:pptx",
    "computer networks slides filetype:pptx",
    "web data mining presentation filetype:pptx",
    "information retrieval lecture filetype:pdf",
    "distributed systems lecture slides pdf",
    "web security presentation pptx",
    "generative models filetype:pptx",
    "diffusion models filetype:pptx ",
    "search engine ranking algorithms lecture filetype:pdf",
    "knowledge graphs presentation filetype:pptx",
    "fake news detection slides filetype:pdf",
    "spam detection web lecture filetype:pptx",
]

# --------------------------------------------------
# STATE (RESUME SAFE)
# --------------------------------------------------
def load_state():
    if STATE_FILE.exists():
        return json.loads(STATE_FILE.read_text())
    return {"seen_urls": [], "slide_count": 0}

def save_state(state):
    STATE_FILE.write_text(json.dumps(state, indent=2))

# --------------------------------------------------
# SERPAPI SEARCH
# --------------------------------------------------
def serpapi_search(query: str) -> List[str]:
    params = {
        "engine": "google",
        "q": query,
        "api_key": SERPAPI_KEY,
        "num": RESULTS_PER_QUERY,
    }
    r = requests.get("https://serpapi.com/search", params=params)
    r.raise_for_status()

    results = r.json().get("organic_results", [])
    links = []

    for r in results:
        link = r.get("link", "")
        if link.lower().endswith((".pptx", ".pdf")):
            links.append(link)

    return links

# --------------------------------------------------
# DOWNLOAD FILE
# --------------------------------------------------
def download_file(url: str) -> Path | None:
    try:
        r = requests.get(url, timeout=20)
        if r.status_code != 200:
            return None

        h = hashlib.md5(url.encode()).hexdigest()[:10]
        ext = url.split(".")[-1].split("?")[0]
        path = OUTPUT_DIR / f"{h}.{ext}"

        with open(path, "wb") as f:
            f.write(r.content)

        return path
    except Exception:
        return None

# --------------------------------------------------
# SLIDE EXTRACTION
# --------------------------------------------------
def extract_pptx_slides(path: Path) -> List[str]:
    prs = Presentation(path)
    slides = []

    for slide in prs.slides:
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
        slides.append("\n".join(text).strip())

    return slides

def extract_pdf_slides(path: Path) -> List[str]:
    text = extract_text(path)
    pages = text.split("\f")
    return [p.strip() for p in pages if p.strip()]

# --------------------------------------------------
# MAIN
# --------------------------------------------------
def main():
    state = load_state()
    print(f"Resuming — slides collected: {state['slide_count']}")

    with OUTPUT_JSONL.open("a", encoding="utf-8") as out:
        for query in QUERIES:
            if state["slide_count"] >= TARGET_SLIDES:
                break

            print(f"\nQuery: {query}")
            links = serpapi_search(query)

            for url in links:
                if state["slide_count"] >= TARGET_SLIDES:
                    break
                if url in state["seen_urls"]:
                    continue

                print(f"Downloading: {url}")
                path = download_file(url)
                if not path:
                    continue

                try:
                    if path.suffix == ".pptx":
                        slides = extract_pptx_slides(path)
                    elif path.suffix == ".pdf":
                        slides = extract_pdf_slides(path)
                    else:
                        continue

                    for i, slide_text in enumerate(slides):
                        if not slide_text.strip():
                            continue
                        if state["slide_count"] >= TARGET_SLIDES:
                            break

                        record = {
                            "source": "google",
                            "method": "serpapi",
                            "query": query,
                            "url": url,
                            "file": str(path),
                            "slide_index": i,
                            "text": slide_text,
                        }

                        out.write(json.dumps(record, ensure_ascii=False) + "\n")
                        out.flush()

                        state["slide_count"] += 1
                        save_state(state)

                    state["seen_urls"].append(url)
                    save_state(state)

                except Exception as e:
                    print(f"Failed extracting {url}: {e}")

                time.sleep(SLEEP)

    print(f"\nDone. Total slides collected: {state['slide_count']}")

if __name__ == "__main__":
    main()
